from pathlib import Path

import pytest
from pptx import Presentation

from tools.sie_autoppt.batch.export import build_shape_map_from_pptx, build_svg_manifest, verify_export_manifest_hash
from tools.sie_autoppt.batch.pptmaster_bridge import (
    REQUIRED_PPTMASTER_SCRIPTS,
    BridgeConfig,
    resolve_bridge_root,
    run_pptmaster_bridge,
)
from tools.sie_autoppt.batch.workspace import BatchWorkspace


def test_resolve_bridge_root_raises_when_root_missing(tmp_path: Path):
    with pytest.raises(FileNotFoundError, match="SIE_PPTMASTER_ROOT"):
        resolve_bridge_root(config=BridgeConfig(pptmaster_root=str(tmp_path / "missing")))


def test_verify_export_manifest_hash_rejects_mismatched_file(tmp_path: Path):
    pptx_path = tmp_path / "exported_raw.pptx"
    pptx_path.write_bytes(b"pptx")
    with pytest.raises(ValueError, match="export hash mismatch"):
        verify_export_manifest_hash(
            run_dir=tmp_path,
            pptx_path="exported_raw.pptx",
            expected_hash="sha256:deadbeef",
        )


def test_build_svg_manifest_hashes_svg_pages(tmp_path: Path):
    svg_dir = tmp_path / "svg_final"
    svg_dir.mkdir()
    (svg_dir / "slide_01.svg").write_text("<svg>1</svg>", encoding="utf-8")
    (svg_dir / "slide_02.svg").write_text("<svg>2</svg>", encoding="utf-8")

    manifest = build_svg_manifest(
        run_id="run-001",
        bundle_hash="sha256:" + ("a" * 64),
        run_dir=tmp_path,
        project_root=tmp_path,
        svg_dir=svg_dir,
        page_refs=["s-001", "s-002"],
    )

    assert manifest.pages[0].page_ref == "s-001"
    assert manifest.pages[1].svg_hash.startswith("sha256:")


def test_build_shape_map_from_pptx_uses_actual_shape_indices(tmp_path: Path):
    pptx_path = tmp_path / "exported_raw.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box1 = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    box1.text_frame.text = "Title"
    box2 = slide.shapes.add_textbox(0, 1000000, 1000000, 1000000)
    box2.text_frame.text = "Body"
    prs.save(pptx_path)

    shape_map = build_shape_map_from_pptx(
        pptx_path=pptx_path,
        page_refs=["s-001"],
    )

    assert len(shape_map) == 2
    assert shape_map[0]["page_ref"] == "s-001"
    assert shape_map[0]["ppt_shape_index"] == 1
    assert shape_map[0]["role"] == "title"
    assert shape_map[1]["ppt_shape_index"] == 2
    assert shape_map[1]["role"] == "body"


def test_run_pptmaster_bridge_supports_compiled_validated_deck(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    run_id = "run-bridge-001"
    workspace = BatchWorkspace.create(root=tmp_path, run_id=run_id)
    bridge_root = tmp_path / "pptmaster"
    for script in REQUIRED_PPTMASTER_SCRIPTS:
        script_path = bridge_root / script
        script_path.parent.mkdir(parents=True, exist_ok=True)
        script_path.write_text("# stub\n", encoding="utf-8")

    def _noop_bridge_command(command: list[str], *, step_name: str) -> None:
        _ = (command, step_name)

    monkeypatch.setattr(
        "tools.sie_autoppt.batch.pptmaster_bridge._run_bridge_command",
        _noop_bridge_command,
    )
    monkeypatch.setattr(
        "tools.sie_autoppt.batch.pptmaster_bridge.sha256_file",
        lambda _path: "sha256:" + ("c" * 64),
    )
    monkeypatch.setattr(
        "tools.sie_autoppt.batch.pptmaster_bridge.build_shape_map_from_pptx",
        lambda **_: [
            {
                "page_ref": "s-001",
                "svg_node_id": "node-1",
                "ppt_shape_name": "Shape 1",
                "ppt_shape_index": 1,
                "role": "title",
            }
        ],
    )

    bundle = {
        "run_id": run_id,
        "bundle_version": 1,
        "bundle_hash": "sha256:" + ("b" * 64),
        "language": "zh-CN",
        "topic": "Bridge smoke",
        "audience": "Engineering",
        "theme": "sie_consulting_fixed",
        "source_index": [
            {
                "source_ref": "src-topic",
                "type": "text",
                "content_hash": "sha256:" + ("d" * 64),
            }
        ],
        "text_summary": {
            "summary": "Bridge smoke summary",
            "key_points": ["Bridge smoke"],
            "source_refs": ["src-topic"],
        },
        "images": [],
        "story_plan": {
            "outline": [
                {
                    "slide_ref": "s-001",
                    "intent": "section",
                    "title": "Bridge smoke",
                    "goal": "Validate bridge",
                    "source_refs": ["src-topic"],
                }
            ]
        },
        "semantic_payload": {
            "meta": {
                "title": "Bridge smoke",
                "theme": "sie_consulting_fixed",
                "language": "zh-CN",
                "author": "AI Auto PPT",
                "version": "2.0",
            },
            "slides": [
                {
                    "slide_id": "s-001",
                    "intent": "section",
                    "title": "Bridge smoke",
                    "subtitle": "Internal batch",
                    "blocks": [],
                }
            ],
        },
    }

    payload = run_pptmaster_bridge(workspace=workspace, bundle=bundle, bridge_root=bridge_root)

    assert payload["pptx_path"] == "bridge/exported_raw.pptx"
    assert payload["svg_manifest"]["pages"][0]["page_ref"] == "s-001"
