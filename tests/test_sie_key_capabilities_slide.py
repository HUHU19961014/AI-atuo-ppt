import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.build_sie_key_capabilities_slide import build_slide, self_check_layout


class SieKeyCapabilitiesSlideTests(unittest.TestCase):
    def test_build_slide_outputs_editable_single_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "key_capabilities.pptx"
            build_slide(output)

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[0].shapes))

    def test_self_check_passes_for_generated_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "key_capabilities.pptx"
            build_slide(output)

            prs = Presentation(output)
            self.assertEqual(self_check_layout(prs), [])


if __name__ == "__main__":
    unittest.main()
