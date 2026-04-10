import tempfile
import unittest
from dataclasses import replace
from pathlib import Path

from pptx import Presentation

from tools.scenario_generators.build_eu_supply_chain_compliance_slide import EU_BRIEF
from tools.scenario_generators.sie_onepage_designer import ACCENT, build_onepage_slide, resolve_reference_policy


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


class SieOnepageDesignerTests(unittest.TestCase):
    def test_reference_policy_requires_explicit_replicate_request(self):
        self.assertEqual(resolve_reference_policy("参考这个风格做一版"), "extract_style_only")
        self.assertEqual(resolve_reference_policy("请保持一致，按这版复刻"), "replicate_allowed")

    def test_title_uses_sie_accent_color(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "accent_title.pptx"
            built, _, _, _ = build_onepage_slide(EU_BRIEF, output_path=output)

            prs = Presentation(built)
            title_runs = []
            for shape in prs.slides[0].shapes:
                if _shape_text(shape) != EU_BRIEF.title:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    title_runs.extend(run for run in paragraph.runs if run.text.strip())
            self.assertTrue(title_runs)
            self.assertEqual(tuple(int(value) for value in title_runs[0].font.color.rgb), ACCENT)

    def test_same_brief_can_render_distinct_variants(self):
        balanced = replace(EU_BRIEF, variant="balanced_dual_panel")
        signal = replace(EU_BRIEF, variant="signal_band")
        summary_board = replace(EU_BRIEF, variant="summary_board")
        comparison = replace(EU_BRIEF, variant="comparison_split")
        with tempfile.TemporaryDirectory() as tmpdir:
            balanced_output = Path(tmpdir) / "balanced.pptx"
            signal_output = Path(tmpdir) / "signal.pptx"
            summary_output = Path(tmpdir) / "summary_board.pptx"
            comparison_output = Path(tmpdir) / "comparison_split.pptx"
            build_onepage_slide(balanced, output_path=balanced_output)
            build_onepage_slide(signal, output_path=signal_output)
            build_onepage_slide(summary_board, output_path=summary_output)
            build_onepage_slide(comparison, output_path=comparison_output)

            balanced_prs = Presentation(balanced_output)
            signal_prs = Presentation(signal_output)
            summary_prs = Presentation(summary_output)
            comparison_prs = Presentation(comparison_output)
            balanced_positions = sorted((shape.left, shape.top, shape.width, shape.height) for shape in balanced_prs.slides[0].shapes)
            signal_positions = sorted((shape.left, shape.top, shape.width, shape.height) for shape in signal_prs.slides[0].shapes)
            summary_positions = sorted((shape.left, shape.top, shape.width, shape.height) for shape in summary_prs.slides[0].shapes)
            comparison_positions = sorted((shape.left, shape.top, shape.width, shape.height) for shape in comparison_prs.slides[0].shapes)
            self.assertNotEqual(balanced_positions, signal_positions)
            self.assertNotEqual(summary_positions, balanced_positions)
            self.assertNotEqual(comparison_positions, signal_positions)


if __name__ == "__main__":
    unittest.main()
