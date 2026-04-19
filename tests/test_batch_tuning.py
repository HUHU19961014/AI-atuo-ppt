from pathlib import Path
import json

import pytest
from pptx import Presentation

from tools.sie_autoppt.batch.contracts import ExportManifest
from tools.sie_autoppt.batch.hashing import sha256_file
from tools.sie_autoppt.batch.hashing import sha256_text
from tools.sie_autoppt.batch.qa_router import route_qa_issues, run_basic_qa, run_pre_export_semantic_qa
from tools.sie_autoppt.batch.tuning import run_deterministic_tuning, verify_manifest_before_tuning
from tools.sie_autoppt.batch.workspace import BatchWorkspace
from tools.sie_autoppt.v2.content_rewriter import RewriteAction, RewriteDeckResult
from tools.sie_autoppt.v2.quality_checks import ContentWarning, QualityGateResult
from tools.sie_autoppt.v2.schema import validate_deck_payload
from tools.sie_autoppt.v2.theme_loader import load_theme


def _write_svg_manifest_for_workspace(
    *,
    workspace: BatchWorkspace,
    run_id: str,
    bundle_hash: str,
    page_refs: tuple[str, ...] = ("s-001",),
) -> str:
    svg_dir = workspace.bridge_dir / "svg_project" / "svg_final"
    svg_dir.mkdir(parents=True, exist_ok=True)
    pages = []
    for index, page_ref in enumerate(page_refs, start=1):
        svg_path = svg_dir / f"slide_{index:02d}.svg"
        svg_path.write_text(f"<svg>{page_ref}</svg>", encoding="utf-8")
        pages.append(
            {
                "page_ref": page_ref,
                "svg_path": svg_path.relative_to(workspace.run_dir).as_posix(),
                "svg_hash": sha256_file(svg_path),
            }
        )
    svg_bundle_hash = sha256_text(json.dumps(pages, ensure_ascii=False, sort_keys=True))
    (workspace.bridge_dir / "svg_manifest.json").write_text(
        json.dumps(
            {
                "run_id": run_id,
                "bundle_hash": bundle_hash,
                "svg_bundle_hash": svg_bundle_hash,
                "project_root": "bridge/svg_project",
                "pages": pages,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return svg_bundle_hash


def test_route_qa_issues_sends_style_issue_to_tune():
    route = route_qa_issues(
        [
            {
                "issue_id": "qa-001",
                "class": "style",
                "severity": "warning",
                "repair_route": "tune",
                "page_ref": "s-001",
                "message": "font mismatch",
            }
        ]
    )
    assert route == "tune"


def test_route_qa_issues_sends_layout_issue_to_regenerate():
    route = route_qa_issues(
        [
            {
                "issue_id": "qa-002",
                "class": "layout",
                "severity": "high",
                "repair_route": "regenerate",
                "page_ref": "s-001",
                "message": "content overflow",
            }
        ]
    )
    assert route == "regenerate"


def test_verify_manifest_before_tuning_rejects_hash_mismatch(tmp_path: Path):
    pptx = tmp_path / "exported_raw.pptx"
    pptx.write_bytes(b"pptx")
    manifest = ExportManifest(
        run_id="run-001",
        bundle_hash="sha256:" + ("a" * 64),
        svg_bundle_hash="sha256:" + ("b" * 64),
        export_hash="sha256:deadbeef",
        exporter_version="pptmaster-bridge-v1",
        pptx_path=str(pptx),
        shape_map=[
            {
                "page_ref": "s-001",
                "svg_node_id": "node-1",
                "ppt_shape_name": "Shape 1",
                "ppt_shape_index": 1,
                "role": "title",
            }
        ],
    )
    with pytest.raises(ValueError, match="export hash mismatch"):
        verify_manifest_before_tuning(manifest, run_dir=tmp_path)


def test_run_deterministic_tuning_normalizes_font(tmp_path: Path):
    pptx_path = tmp_path / "exported_raw.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "hello"
    run.font.name = "Arial"
    prs.save(pptx_path)

    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-001")
    bundle_hash = "sha256:" + ("a" * 64)
    svg_bundle_hash = _write_svg_manifest_for_workspace(
        workspace=workspace,
        run_id="run-001",
        bundle_hash=bundle_hash,
    )
    result = run_deterministic_tuning(
        workspace=workspace,
        export_manifest={
            "run_id": "run-001",
            "bundle_hash": bundle_hash,
            "svg_bundle_hash": svg_bundle_hash,
            "export_hash": sha256_file(pptx_path),
            "exporter_version": "pptmaster-bridge-v1",
            "pptx_path": str(pptx_path),
            "shape_map": [
                {
                    "page_ref": "s-001",
                    "svg_node_id": "n1",
                    "ppt_shape_name": "TextBox 1",
                    "ppt_shape_index": 1,
                    "role": "title",
                }
            ],
        },
    )

    tuned = Presentation((workspace.run_dir / result["pptx_path"]).resolve())
    text_runs = []
    for shape in tuned.slides[0].shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text_runs.extend(paragraph.runs)
    assert any(run.font.name == "Microsoft YaHei" for run in text_runs)


def test_run_deterministic_tuning_applies_role_colors(tmp_path: Path):
    theme = load_theme("sie_consulting_fixed")
    pptx_path = tmp_path / "exported_raw_colored.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "Title"
    body_box = slide.shapes.add_textbox(0, 1000000, 1000000, 1000000)
    body_run = body_box.text_frame.paragraphs[0].add_run()
    body_run.text = "Body"
    prs.save(pptx_path)

    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-colors")
    bundle_hash = "sha256:" + ("a" * 64)
    svg_bundle_hash = _write_svg_manifest_for_workspace(
        workspace=workspace,
        run_id="run-colors",
        bundle_hash=bundle_hash,
    )
    result = run_deterministic_tuning(
        workspace=workspace,
        export_manifest={
            "run_id": "run-colors",
            "bundle_hash": bundle_hash,
            "svg_bundle_hash": svg_bundle_hash,
            "export_hash": sha256_file(pptx_path),
            "exporter_version": "pptmaster-bridge-v1",
            "pptx_path": str(pptx_path),
            "shape_map": [
                {
                    "page_ref": "s-001",
                    "svg_node_id": "n1",
                    "ppt_shape_name": "TextBox 1",
                    "ppt_shape_index": 1,
                    "role": "title",
                },
                {
                    "page_ref": "s-001",
                    "svg_node_id": "n2",
                    "ppt_shape_name": "TextBox 2",
                    "ppt_shape_index": 2,
                    "role": "body",
                },
            ],
        },
    )

    tuned = Presentation((workspace.run_dir / result["pptx_path"]).resolve())
    title_shape = tuned.slides[0].shapes[0]
    body_shape = tuned.slides[0].shapes[1]
    title_color = str(title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb)
    body_color = str(body_shape.text_frame.paragraphs[0].runs[0].font.color.rgb)
    assert title_color == theme.colors.primary.removeprefix("#")
    assert body_color == theme.colors.text_main.removeprefix("#")


def test_run_basic_qa_routes_font_mismatch_to_tune(tmp_path: Path):
    pptx_path = tmp_path / "raw.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "hello"
    run.font.name = "Arial"
    prs.save(pptx_path)

    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-qa-style")
    result = run_basic_qa(
        workspace=workspace,
        export_manifest={
            "pptx_path": str(pptx_path),
            "shape_map": [
                {
                    "page_ref": "s-001",
                    "svg_node_id": "n1",
                    "ppt_shape_name": "TextBox 1",
                    "ppt_shape_index": 1,
                    "role": "title",
                }
            ],
        },
        tuning_result={"pptx_path": str(pptx_path)},
    )

    assert result["status"] == "repairable"
    assert result["route"] == "tune"
    assert result["issues"][0]["class"] == "style"


def test_run_basic_qa_routes_missing_shape_map_to_regenerate(tmp_path: Path):
    pptx_path = tmp_path / "raw.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(pptx_path)

    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-qa-layout")
    result = run_basic_qa(
        workspace=workspace,
        export_manifest={
            "pptx_path": str(pptx_path),
            "shape_map": [],
        },
        tuning_result={"pptx_path": str(pptx_path)},
    )

    assert result["status"] == "repairable"
    assert result["route"] == "regenerate"
    assert result["issues"][0]["repair_route"] == "regenerate"


def test_run_basic_qa_reports_degraded_mode_when_shape_map_is_heuristic(tmp_path: Path):
    pptx_path = tmp_path / "raw.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(pptx_path)

    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-qa-heuristic")
    result = run_basic_qa(
        workspace=workspace,
        export_manifest={
            "pptx_path": str(pptx_path),
            "shape_map": [
                {
                    "page_ref": "s-001",
                    "svg_node_id": "n1",
                    "ppt_shape_name": "TextBox 1",
                    "ppt_shape_index": 1,
                    "role": "title",
                }
            ],
            "shape_map_mode": "heuristic",
        },
        tuning_result={"pptx_path": str(pptx_path)},
    )

    assert result["degraded_mode"] is True
    assert result["degraded_reasons"]
    assert "heuristic" in result["degraded_reasons"][0]


def test_verify_manifest_before_tuning_rejects_svg_bundle_hash_mismatch(tmp_path: Path):
    workspace = BatchWorkspace.create(root=tmp_path, run_id="run-svg-hash")
    pptx_path = workspace.bridge_dir / "exported_raw.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(pptx_path)

    (workspace.bridge_dir / "svg_manifest.json").write_text(
        json.dumps(
            {
                "run_id": "run-svg-hash",
                "bundle_hash": "sha256:" + ("a" * 64),
                "svg_bundle_hash": "sha256:" + ("1" * 64),
                "project_root": "bridge/svg_project",
                "pages": [
                    {
                        "page_ref": "s-001",
                        "svg_path": "bridge/svg_project/svg_final/slide_01.svg",
                        "svg_hash": "sha256:" + ("2" * 64),
                    }
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    manifest = ExportManifest(
        run_id="run-svg-hash",
        bundle_hash="sha256:" + ("a" * 64),
        svg_bundle_hash="sha256:" + ("f" * 64),
        export_hash=sha256_file(pptx_path),
        exporter_version="pptmaster-bridge-v1",
        pptx_path="bridge/exported_raw.pptx",
        shape_map=[
            {
                "page_ref": "s-001",
                "svg_node_id": "node-1",
                "ppt_shape_name": "Shape 1",
                "ppt_shape_index": 1,
                "role": "title",
            }
        ],
    )

    with pytest.raises(ValueError, match="svg bundle hash mismatch"):
        verify_manifest_before_tuning(manifest, run_dir=workspace.run_dir)


def _sample_validated_deck():
    return validate_deck_payload(
        {
            "meta": {
                "title": "QA Deck",
                "theme": "sie_consulting_fixed",
                "language": "zh-CN",
                "author": "AI Auto PPT",
                "version": "2.0",
            },
            "slides": [{"slide_id": "s-001", "layout": "title_only", "title": "项目建议"}],
        }
    )


def _quality_result(*, validated_deck, passed: bool, warnings=(), high=(), errors=()):
    return QualityGateResult(
        passed=passed,
        review_required=not passed,
        warnings=tuple(warnings),
        high=tuple(high),
        errors=tuple(errors),
        validated_deck=validated_deck,
    )


def test_run_pre_export_semantic_qa_converges_within_rewrite_round_limit(monkeypatch: pytest.MonkeyPatch):
    validated = _sample_validated_deck()
    initial_quality = _quality_result(
        validated_deck=validated,
        passed=False,
        high=(ContentWarning(slide_id="s-001", warning_level="high", message="content overflow"),),
    )
    round1_quality = _quality_result(
        validated_deck=validated,
        passed=False,
        warnings=(ContentWarning(slide_id="s-001", warning_level="warning", message="title too generic"),),
    )
    round2_quality = _quality_result(
        validated_deck=validated,
        passed=True,
    )

    rewrite_results = [
        RewriteDeckResult(
            attempted=True,
            applied=True,
            validated_deck=validated,
            initial_quality_gate=initial_quality,
            final_quality_gate=round1_quality,
            actions=(
                RewriteAction(
                    slide_id="s-001",
                    field="title",
                    action="rewrite_title",
                    before="项目背景",
                    after="项目建议",
                ),
            ),
            notes=("rewrite pass completed",),
        ),
        RewriteDeckResult(
            attempted=True,
            applied=True,
            validated_deck=validated,
            initial_quality_gate=round1_quality,
            final_quality_gate=round2_quality,
            actions=(
                RewriteAction(
                    slide_id="s-001",
                    field="title",
                    action="refine_title",
                    before="项目建议",
                    after="项目建议（已收敛）",
                ),
            ),
            notes=("rewrite pass completed",),
        ),
    ]

    monkeypatch.setattr(
        "tools.sie_autoppt.batch.qa_router.compile_semantic_deck_payload",
        lambda *_args, **_kwargs: type("CompiledDeck", (), {"deck": validated.deck})(),
    )
    monkeypatch.setattr("tools.sie_autoppt.batch.qa_router.quality_gate", lambda _deck: initial_quality)
    monkeypatch.setattr(
        "tools.sie_autoppt.batch.qa_router.rewrite_deck",
        lambda _deck, _quality: rewrite_results.pop(0),
    )

    report = run_pre_export_semantic_qa(
        run_id="run-pre-qa-pass",
        bundle={
            "semantic_payload": {"meta": {"title": "QA Deck", "theme": "sie_consulting_fixed", "language": "zh-CN"}},
            "topic": "QA Deck",
            "theme": "sie_consulting_fixed",
            "language": "zh-CN",
        },
        max_rewrite_rounds=2,
    )

    assert report["status"] == "passed"
    assert report["route"] == "stop"
    assert report["rewrite_round_limit"] == 2
    assert report["rewrite_rounds_used"] == 2
    assert len(report["rewrite_rounds"]) == 2
    assert report["rewrite_rounds"][0]["action_count"] == 1
    assert report["rewrite_rounds"][1]["action_count"] == 1


def test_run_pre_export_semantic_qa_stops_fast_on_unfixable_errors(monkeypatch: pytest.MonkeyPatch):
    validated = _sample_validated_deck()
    initial_quality = _quality_result(
        validated_deck=validated,
        passed=False,
        errors=(ContentWarning(slide_id="s-001", warning_level="error", message="schema mismatch"),),
    )
    stuck_quality = _quality_result(
        validated_deck=validated,
        passed=False,
        errors=(ContentWarning(slide_id="s-001", warning_level="error", message="schema mismatch"),),
    )

    rewrite_calls = {"count": 0}

    def _fake_rewrite(_deck, _quality):
        rewrite_calls["count"] += 1
        return RewriteDeckResult(
            attempted=True,
            applied=False,
            validated_deck=validated,
            initial_quality_gate=initial_quality,
            final_quality_gate=stuck_quality,
            actions=(),
            notes=("fixable issues were detected but no safe rewrite was produced",),
        )

    monkeypatch.setattr(
        "tools.sie_autoppt.batch.qa_router.compile_semantic_deck_payload",
        lambda *_args, **_kwargs: type("CompiledDeck", (), {"deck": validated.deck})(),
    )
    monkeypatch.setattr("tools.sie_autoppt.batch.qa_router.quality_gate", lambda _deck: initial_quality)
    monkeypatch.setattr("tools.sie_autoppt.batch.qa_router.rewrite_deck", _fake_rewrite)

    report = run_pre_export_semantic_qa(
        run_id="run-pre-qa-fail",
        bundle={
            "semantic_payload": {"meta": {"title": "QA Deck", "theme": "sie_consulting_fixed", "language": "zh-CN"}},
            "topic": "QA Deck",
            "theme": "sie_consulting_fixed",
            "language": "zh-CN",
        },
        max_rewrite_rounds=2,
    )

    assert report["status"] == "failed"
    assert report["route"] == "stop"
    assert report["rewrite_round_limit"] == 2
    assert report["rewrite_rounds_used"] == 1
    assert len(report["rewrite_rounds"]) == 1
    assert rewrite_calls["count"] == 1
