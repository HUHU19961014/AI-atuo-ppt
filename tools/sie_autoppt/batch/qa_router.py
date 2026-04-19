from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from ..v2 import compile_semantic_deck_payload
from ..v2.content_rewriter import RewriteDeckResult, rewrite_deck
from ..v2.quality_checks import QualityGateResult, quality_gate
from .contracts import QaIssue, QaReport
from .export import resolve_run_artifact_path
from .workspace import BatchWorkspace

DEFAULT_PRE_EXPORT_REWRITE_MAX_ROUNDS = 2


def route_qa_issues(issues: list[dict[str, str]]) -> str:
    if any(issue.get("repair_route") == "stop" for issue in issues):
        return "stop"
    if any(issue.get("repair_route") == "regenerate" for issue in issues):
        return "regenerate"
    return "tune"


def run_pre_export_semantic_qa(
    *,
    run_id: str,
    bundle: dict[str, Any],
    max_rewrite_rounds: int = DEFAULT_PRE_EXPORT_REWRITE_MAX_ROUNDS,
) -> dict[str, Any]:
    rewrite_round_limit = max(0, int(max_rewrite_rounds))
    rewrite_rounds: list[dict[str, Any]] = []

    semantic_payload = dict(bundle.get("semantic_payload") or {})
    meta = semantic_payload.get("meta") or {}
    deck = compile_semantic_deck_payload(
        semantic_payload,
        default_title=str(bundle.get("topic") or meta.get("title") or "Untitled"),
        default_theme=str(bundle.get("theme") or meta.get("theme") or "sie_consulting_fixed"),
        default_language=str(bundle.get("language") or meta.get("language") or "zh-CN"),
        default_author=str(meta.get("author") or "AI Auto PPT"),
    ).deck

    current_quality = _run_pre_export_quality_rules(deck)
    if current_quality.passed:
        return _report_for_issues(
            run_id=run_id,
            issues=[],
            rewrite_round_limit=rewrite_round_limit,
            rewrite_rounds=rewrite_rounds,
        ).model_dump(mode="json", by_alias=True)

    current_validated_deck = current_quality.validated_deck
    if current_validated_deck is not None and rewrite_round_limit > 0:
        for round_index in range(1, rewrite_round_limit + 1):
            rewrite_result = _execute_pre_export_rewrite_round(current_validated_deck, current_quality)
            rewrite_rounds.append(_rewrite_round_report(round_index=round_index, rewrite_result=rewrite_result))
            current_quality = rewrite_result.final_quality_gate
            current_validated_deck = rewrite_result.validated_deck

            if current_quality.passed:
                return _report_for_issues(
                    run_id=run_id,
                    issues=[],
                    rewrite_round_limit=rewrite_round_limit,
                    rewrite_rounds=rewrite_rounds,
                ).model_dump(mode="json", by_alias=True)
            if not rewrite_result.applied or current_validated_deck is None:
                break

    issues = _quality_issues(current_quality)
    return _report_for_issues(
        run_id=run_id,
        issues=issues,
        rewrite_round_limit=rewrite_round_limit,
        rewrite_rounds=rewrite_rounds,
    ).model_dump(mode="json", by_alias=True)


def run_basic_qa(
    *,
    workspace: BatchWorkspace,
    export_manifest: dict[str, Any],
    tuning_result: dict[str, Any] | Any,
) -> dict[str, Any]:
    candidate = str(export_manifest.get("pptx_path") or "")
    if isinstance(tuning_result, dict) and tuning_result.get("pptx_path"):
        candidate = str(tuning_result["pptx_path"])
    pptx_path = resolve_run_artifact_path(run_dir=workspace.run_dir, artifact_path=candidate)

    issues: list[QaIssue] = []
    shape_map_mode = str(export_manifest.get("shape_map_mode", "heuristic") or "heuristic").strip().lower()
    if shape_map_mode not in {"mapped", "heuristic"}:
        shape_map_mode = "heuristic"
    degraded_reasons: list[str] = []
    if shape_map_mode == "heuristic":
        degraded_reasons.append("shape_map_mode=heuristic; fallback shape mapping in use.")
    if not pptx_path.exists():
        issues.append(
            QaIssue(
                issue_id="qa-001",
                class_="schema",
                severity="error",
                repair_route="stop",
                message=f"missing PPTX artifact: {pptx_path}",
            )
        )
    if not export_manifest.get("shape_map"):
        issues.append(
            QaIssue(
                issue_id=f"qa-{len(issues) + 1:03d}",
                class_="mapping",
                severity="high",
                repair_route="regenerate",
                message="shape_map is empty",
            )
        )
    if pptx_path.exists():
        issues.extend(_collect_font_issues(pptx_path, start_index=len(issues) + 1))

    report = _report_for_issues(
        run_id=str(export_manifest.get("run_id") or workspace.run_dir.name),
        issues=issues,
        degraded_reasons=degraded_reasons,
    )
    report.checked_pptx_path = candidate
    report.workspace = workspace.run_dir.as_posix()
    return report.model_dump(mode="json", by_alias=True)


def _collect_font_issues(pptx_path: Path, *, start_index: int) -> list[QaIssue]:
    presentation = Presentation(pptx_path)
    issues: list[QaIssue] = []
    issue_index = start_index
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_name = run.font.name or ""
                    if font_name and font_name == "Microsoft YaHei":
                        continue
                    if not run.text.strip():
                        continue
                    issues.append(
                        QaIssue(
                            issue_id=f"qa-{issue_index:03d}",
                            class_="style",
                            severity="warning",
                            repair_route="tune",
                            page_ref=f"s-{slide_index:03d}",
                            message=f"font mismatch on shape {shape_index}: {font_name or 'unset'}",
                        )
                    )
                    issue_index += 1
    return issues


def _run_pre_export_quality_rules(deck) -> QualityGateResult:
    return quality_gate(deck)


def _execute_pre_export_rewrite_round(validated_deck, quality_result: QualityGateResult) -> RewriteDeckResult:
    return rewrite_deck(validated_deck, quality_result)


def _quality_issues(result: QualityGateResult) -> list[QaIssue]:
    issues: list[QaIssue] = []
    counter = 1
    for warning in result.warnings:
        issues.append(
            QaIssue(
                issue_id=f"qa-pre-{counter:03d}",
                class_="content",
                severity="warning",
                repair_route="tune",
                page_ref=warning.slide_id,
                message=warning.message,
            )
        )
        counter += 1
    for warning in result.high:
        issues.append(
            QaIssue(
                issue_id=f"qa-pre-{counter:03d}",
                class_="content",
                severity="high",
                repair_route="regenerate",
                page_ref=warning.slide_id,
                message=warning.message,
            )
        )
        counter += 1
    for warning in result.errors:
        issues.append(
            QaIssue(
                issue_id=f"qa-pre-{counter:03d}",
                class_="content",
                severity="error",
                repair_route="stop",
                page_ref=warning.slide_id,
                message=warning.message,
            )
        )
        counter += 1
    return issues


def _rewrite_round_report(*, round_index: int, rewrite_result: RewriteDeckResult) -> dict[str, Any]:
    return {
        "round": round_index,
        "attempted": rewrite_result.attempted,
        "applied": rewrite_result.applied,
        "action_count": len(rewrite_result.actions),
        "rewritten_slide_ids": sorted({action.slide_id for action in rewrite_result.actions}),
        "initial_summary": rewrite_result.initial_quality_gate.summary,
        "final_summary": rewrite_result.final_quality_gate.summary,
        "notes": list(rewrite_result.notes),
    }


def _report_for_issues(
    *,
    run_id: str,
    issues: list[QaIssue],
    degraded_reasons: list[str] | None = None,
    rewrite_round_limit: int = 0,
    rewrite_rounds: list[dict[str, Any]] | None = None,
) -> QaReport:
    normalized_reasons = [str(reason).strip() for reason in (degraded_reasons or []) if str(reason).strip()]
    degraded_mode = bool(normalized_reasons)
    normalized_rewrite_rounds = list(rewrite_rounds or [])
    rewrite_rounds_used = len(normalized_rewrite_rounds)
    if not issues:
        return QaReport(
            run_id=run_id,
            status="passed",
            route="stop",
            issues=[],
            degraded_mode=degraded_mode,
            degraded_reasons=normalized_reasons,
            rewrite_round_limit=max(0, int(rewrite_round_limit)),
            rewrite_rounds_used=rewrite_rounds_used,
            rewrite_rounds=normalized_rewrite_rounds,
        )
    payloads = [issue.model_dump(mode="json", by_alias=True) for issue in issues]
    route = route_qa_issues(payloads)
    status = "failed" if route == "stop" else "repairable"
    return QaReport(
        run_id=run_id,
        status=status,
        route=route,
        issues=issues,
        degraded_mode=degraded_mode,
        degraded_reasons=normalized_reasons,
        rewrite_round_limit=max(0, int(rewrite_round_limit)),
        rewrite_rounds_used=rewrite_rounds_used,
        rewrite_rounds=normalized_rewrite_rounds,
    )
