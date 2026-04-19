from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from .contracts import SvgManifest
from .hashing import sha256_file, sha256_text


def resolve_run_artifact_path(*, run_dir: Path, artifact_path: str) -> Path:
    candidate = Path(artifact_path)
    return candidate if candidate.is_absolute() else run_dir / candidate


def verify_export_manifest_hash(*, run_dir: Path, pptx_path: str, expected_hash: str) -> None:
    actual_path = resolve_run_artifact_path(run_dir=run_dir, artifact_path=pptx_path)
    actual = sha256_file(actual_path)
    if actual != expected_hash:
        raise ValueError(f"export hash mismatch: expected {expected_hash}, got {actual}")


def build_svg_manifest(
    *,
    run_id: str,
    bundle_hash: str,
    run_dir: Path,
    project_root: Path,
    svg_dir: Path,
    page_refs: list[str],
) -> SvgManifest:
    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        raise ValueError("svg output is empty")
    pages: list[dict[str, str]] = []
    for index, svg_path in enumerate(svg_files):
        page_ref = page_refs[index] if index < len(page_refs) else f"s-{index + 1:03d}"
        pages.append(
            {
                "page_ref": page_ref,
                "svg_path": svg_path.relative_to(run_dir).as_posix(),
                "svg_hash": sha256_file(svg_path),
            }
        )
    svg_bundle_hash = sha256_text(json.dumps(pages, ensure_ascii=False, sort_keys=True))
    return SvgManifest(
        run_id=run_id,
        bundle_hash=bundle_hash,
        svg_bundle_hash=svg_bundle_hash,
        project_root=project_root.relative_to(run_dir).as_posix(),
        pages=pages,
    )


def build_shape_map_from_pptx(*, pptx_path: Path, page_refs: list[str]) -> list[dict[str, str | int]]:
    presentation = Presentation(pptx_path)
    shape_map: list[dict[str, str | int]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        page_ref = page_refs[slide_index - 1] if slide_index - 1 < len(page_refs) else f"s-{slide_index:03d}"
        seen_text_shape = False
        for shape_index, shape in enumerate(slide.shapes, start=1):
            role = "graphic"
            if getattr(shape, "has_text_frame", False):
                has_text = any(run.text.strip() for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
                if has_text:
                    role = "title" if not seen_text_shape else "body"
                    seen_text_shape = True
            shape_map.append(
                {
                    "page_ref": page_ref,
                    "svg_node_id": f"ppt-slide{slide_index}-shape{shape_index}",
                    "ppt_shape_name": str(getattr(shape, "name", f"Shape {shape_index}")),
                    "ppt_shape_index": shape_index,
                    "role": role,
                }
            )
    return shape_map
