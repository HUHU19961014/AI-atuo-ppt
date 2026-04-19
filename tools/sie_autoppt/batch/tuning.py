from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor

from ..v2.theme_loader import load_theme
from .contracts import ExportManifest, SvgManifest
from .export import resolve_run_artifact_path
from .hashing import sha256_file, sha256_text
from .workspace import BatchWorkspace


def verify_manifest_before_tuning(manifest: ExportManifest, *, run_dir: Path) -> None:
    pptx_path = resolve_run_artifact_path(run_dir=run_dir, artifact_path=manifest.pptx_path)
    actual_hash = sha256_file(pptx_path)
    if actual_hash != manifest.export_hash:
        raise ValueError(f"export hash mismatch: expected {manifest.export_hash}, got {actual_hash}")
    _verify_svg_bundle_before_tuning(manifest=manifest, run_dir=run_dir)


def run_deterministic_tuning(*, workspace: BatchWorkspace, export_manifest: dict[str, Any]) -> dict[str, Any]:
    manifest = ExportManifest.model_validate(export_manifest)
    verify_manifest_before_tuning(manifest, run_dir=workspace.run_dir)

    source = resolve_run_artifact_path(run_dir=workspace.run_dir, artifact_path=manifest.pptx_path)
    target = workspace.tune_dir / "tuned.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, target)

    presentation = Presentation(target)
    theme = load_theme("sie_consulting_fixed")
    actions: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name != "Microsoft YaHei":
                        run.font.name = "Microsoft YaHei"
                        actions.append("font_name_normalized")
    actions.extend(_apply_role_styles(presentation, manifest, theme))
    presentation.save(target)

    return {
        "status": "tuned",
        "actions": actions,
        "pptx_path": target.relative_to(workspace.run_dir).as_posix(),
        "tuned_hash": sha256_file(target),
    }


def run_noop_tuning(*, workspace: BatchWorkspace, export_manifest: dict[str, Any]) -> dict[str, Any]:
    return run_deterministic_tuning(workspace=workspace, export_manifest=export_manifest)


def _verify_svg_bundle_before_tuning(*, manifest: ExportManifest, run_dir: Path) -> None:
    svg_manifest_path = run_dir / "bridge" / "svg_manifest.json"
    if not svg_manifest_path.exists():
        raise ValueError(f"svg manifest missing for tuning verification: {svg_manifest_path}")
    svg_manifest = SvgManifest.model_validate_json(svg_manifest_path.read_text(encoding="utf-8"))
    if svg_manifest.svg_bundle_hash != manifest.svg_bundle_hash:
        raise ValueError(
            f"svg bundle hash mismatch: expected {manifest.svg_bundle_hash}, got {svg_manifest.svg_bundle_hash}"
        )
    if svg_manifest.bundle_hash != manifest.bundle_hash:
        raise ValueError(f"bundle hash mismatch: expected {manifest.bundle_hash}, got {svg_manifest.bundle_hash}")

    normalized_pages: list[dict[str, str]] = []
    for page in svg_manifest.pages:
        page_path = resolve_run_artifact_path(run_dir=run_dir, artifact_path=page.svg_path)
        if not page_path.exists():
            raise ValueError(f"svg manifest page missing: {page_path}")
        actual_svg_hash = sha256_file(page_path)
        if actual_svg_hash != page.svg_hash:
            raise ValueError(
                f"svg page hash mismatch for {page.svg_path}: expected {page.svg_hash}, got {actual_svg_hash}"
            )
        normalized_pages.append(
            {
                "page_ref": page.page_ref,
                "svg_path": page.svg_path,
                "svg_hash": page.svg_hash,
            }
        )

    actual_svg_bundle_hash = sha256_text(json.dumps(normalized_pages, ensure_ascii=False, sort_keys=True))
    if actual_svg_bundle_hash != manifest.svg_bundle_hash:
        raise ValueError(
            f"svg bundle hash mismatch: expected {manifest.svg_bundle_hash}, got {actual_svg_bundle_hash}"
        )


def _apply_role_styles(presentation: Presentation, manifest: ExportManifest, theme) -> list[str]:
    actions: list[str] = []
    page_refs_in_order: list[str] = []
    for entry in manifest.shape_map:
        if entry.page_ref not in page_refs_in_order:
            page_refs_in_order.append(entry.page_ref)

    grouped_entries: dict[str, list[Any]] = {}
    for entry in manifest.shape_map:
        grouped_entries.setdefault(entry.page_ref, []).append(entry)

    title_color = RGBColor.from_string(theme.colors.primary.removeprefix("#"))
    body_color = RGBColor.from_string(theme.colors.text_main.removeprefix("#"))

    for slide_index, slide in enumerate(presentation.slides):
        if slide_index >= len(page_refs_in_order):
            continue
        page_ref = page_refs_in_order[slide_index]
        for entry in grouped_entries.get(page_ref, []):
            shape_pos = entry.ppt_shape_index - 1
            if shape_pos < 0 or shape_pos >= len(slide.shapes):
                continue
            shape = slide.shapes[shape_pos]
            if not getattr(shape, "has_text_frame", False):
                continue
            target_color = title_color if entry.role == "title" else body_color
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    run.font.color.rgb = target_color
                    if entry.role == "title":
                        run.font.bold = True
                    actions.append(f"role_style_applied:{entry.role}")
    return actions
