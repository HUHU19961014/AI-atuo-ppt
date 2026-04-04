import time
import zipfile
from copy import deepcopy
from pathlib import Path
from xml.etree import ElementTree

from pptx import Presentation


def clone_slide_after(prs: Presentation, source_idx: int, insert_after_idx: int, keep_rel_ids: bool = True):
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        element = shape._element
        element.getparent().remove(element)
    for shape in source.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            if keep_rel_ids:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            else:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target)
        except Exception:
            pass
    slide_id_list = prs.slides._sldIdLst
    new_id = slide_id_list[-1]
    del slide_id_list[-1]
    slide_id_list.insert(insert_after_idx + 1, new_id)
    return prs.slides[insert_after_idx + 1]


def remove_slide(prs: Presentation, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    prs.part.drop_rel(slide_id.rId)
    del slide_id_list[slide_index]


def ensure_last_slide(prs: Presentation, slide_id: int):
    slide_id_list = prs.slides._sldIdLst
    target = None
    for item in slide_id_list:
        if int(item.id) == int(slide_id):
            target = item
            break
    if target is None:
        return
    slide_id_list.remove(target)
    slide_id_list.append(target)


def slide_image_targets(pptx_path: Path, slide_no: int) -> set[str]:
    rel_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    targets: set[str] = set()
    with zipfile.ZipFile(pptx_path) as package:
        try:
            root = ElementTree.fromstring(package.read(rel_path))
        except KeyError:
            return targets
    for rel in root:
        rel_type = rel.attrib.get("Type", "")
        if rel_type.endswith("/image"):
            target = rel.attrib.get("Target", "")
            if target:
                targets.add(target)
    return targets


def slide_assets_preserved(pptx_path: Path, source_idx: int, target_indices: list[int]) -> bool:
    source_targets = slide_image_targets(pptx_path, source_idx)
    if not source_targets:
        return True
    for slide_no in target_indices:
        target_assets = slide_image_targets(pptx_path, slide_no)
        if not source_targets.issubset(target_assets):
            return False
    return True


def copy_slide_xml_assets(pptx_path: Path, source_idx: int, target_indices: list[int]) -> bool:
    if not target_indices:
        return True

    source_slide_name = f"ppt/slides/slide{source_idx}.xml"
    source_rel_name = f"ppt/slides/_rels/slide{source_idx}.xml.rels"
    target_slide_names = {f"ppt/slides/slide{target}.xml" for target in target_indices}
    target_rel_names = {f"ppt/slides/_rels/slide{target}.xml.rels" for target in target_indices}
    rebuilt_path = pptx_path.with_name(pptx_path.stem + "_rebuilt.pptx")

    with zipfile.ZipFile(pptx_path, "r") as source_package:
        if source_slide_name not in source_package.namelist():
            return False

        slide_root = ElementTree.fromstring(source_package.read(source_slide_name))
        slide_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        source_sp_tree = slide_root.find(f".//{slide_ns}spTree")
        source_pics = [deepcopy(pic) for pic in source_sp_tree.findall(f"{slide_ns}pic")] if source_sp_tree is not None else []
        rel_bytes = source_package.read(source_rel_name) if source_rel_name in source_package.namelist() else None
        source_image_rels = []
        if rel_bytes is not None:
            source_rel_root = ElementTree.fromstring(rel_bytes)
            source_image_rels = [
                deepcopy(rel)
                for rel in source_rel_root
                if rel.attrib.get("Type", "").endswith("/image")
            ]

        slide_replacements: dict[str, bytes] = {}
        rel_replacements: dict[str, bytes] = {}
        rel_root_tag = (
            ElementTree.fromstring(rel_bytes).tag
            if rel_bytes is not None
            else "{http://schemas.openxmlformats.org/package/2006/relationships}Relationships"
        )

        for target_slide_name in target_slide_names:
            if target_slide_name not in source_package.namelist():
                continue
            target_root = ElementTree.fromstring(source_package.read(target_slide_name))
            target_sp_tree = target_root.find(f".//{slide_ns}spTree")
            if target_sp_tree is not None and not target_sp_tree.findall(f"{slide_ns}pic") and source_pics:
                insert_at = next(
                    (index for index, child in enumerate(list(target_sp_tree)) if child.tag == f"{slide_ns}extLst"),
                    len(target_sp_tree),
                )
                for pic in source_pics:
                    target_sp_tree.insert(insert_at, deepcopy(pic))
                    insert_at += 1
            slide_replacements[target_slide_name] = ElementTree.tostring(
                target_root,
                encoding="utf-8",
                xml_declaration=True,
            )

        for target_rel_name in target_rel_names:
            if target_rel_name in source_package.namelist():
                target_rel_root = ElementTree.fromstring(source_package.read(target_rel_name))
            else:
                target_rel_root = ElementTree.Element(rel_root_tag)
            existing_image_targets = {
                rel.attrib.get("Target", "")
                for rel in target_rel_root
                if rel.attrib.get("Type", "").endswith("/image")
            }
            for image_rel in source_image_rels:
                if image_rel.attrib.get("Target", "") in existing_image_targets:
                    continue
                target_rel_root.append(deepcopy(image_rel))
            rel_replacements[target_rel_name] = ElementTree.tostring(
                target_rel_root,
                encoding="utf-8",
                xml_declaration=True,
            )

        with zipfile.ZipFile(rebuilt_path, "w", zipfile.ZIP_DEFLATED) as rebuilt:
            for info in source_package.infolist():
                data = source_package.read(info.filename)
                if info.filename in slide_replacements:
                    data = slide_replacements[info.filename]
                elif info.filename in rel_replacements:
                    data = rel_replacements[info.filename]
                rebuilt.writestr(info, data)

            for target_slide_name, data in slide_replacements.items():
                if target_slide_name not in source_package.namelist():
                    rebuilt.writestr(target_slide_name, data)
            for target_rel_name, data in rel_replacements.items():
                if target_rel_name not in source_package.namelist():
                    rebuilt.writestr(target_rel_name, data)

    last_error = None
    for _ in range(10):
        try:
            rebuilt_path.replace(pptx_path)
            return True
        except PermissionError as exc:
            last_error = exc
            time.sleep(0.5)
    if rebuilt_path.exists():
        rebuilt_path.unlink(missing_ok=True)
    if last_error is not None:
        raise last_error
    return True
