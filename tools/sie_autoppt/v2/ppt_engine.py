from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .io import RenderLog
from .layout_router import render_slide
from .quality_checks import ContentWarning, check_deck_content, count_errors
from .schema import DeckDocument, ValidatedDeck, validate_deck_payload
from .theme_loader import load_theme


@dataclass(frozen=True)
class RenderArtifacts:
    output_path: Path
    log_path: Path | None
    warnings_path: Path | None
    slide_count: int
    warnings: tuple[str, ...] = ()
    content_warnings: tuple[ContentWarning, ...] = ()
    error_count: int = 0


def _coerce_validated_deck(deck_data: DeckDocument | ValidatedDeck | dict[str, object]) -> ValidatedDeck:
    if isinstance(deck_data, ValidatedDeck):
        return deck_data
    if isinstance(deck_data, DeckDocument):
        return ValidatedDeck(deck=deck_data)
    if isinstance(deck_data, dict):
        return validate_deck_payload(deck_data)
    raise TypeError("deck_data must be a DeckDocument, ValidatedDeck, or JSON object.")


def generate_ppt(
    deck_data: DeckDocument | ValidatedDeck | dict[str, object],
    output_path: str | Path,
    theme_name: str | None = None,
    log_path: str | Path | None = None,
    max_errors: int = 0,
) -> RenderArtifacts:
    """
    Generate a PowerPoint presentation from a deck specification.

    Args:
        deck_data: Deck specification (DeckDocument, ValidatedDeck, or dict)
        output_path: Path where the PPTX file will be saved
        theme_name: Optional theme name override
        log_path: Optional path to write the generation log
        max_errors: Maximum number of error-level quality issues allowed (default: 0)
                   If error count exceeds this threshold, generation will be aborted.

    Returns:
        RenderArtifacts containing output paths and quality warnings

    Raises:
        ValueError: If error count exceeds max_errors threshold
    """
    validated = _coerce_validated_deck(deck_data)
    deck = validated.deck
    theme = load_theme(theme_name or deck.meta.theme)

    log = RenderLog()
    log.info(f"deck title: {deck.meta.title}")
    log.info(f"theme: {theme.theme_name}")
    log.extend(validated.warnings)

    # Run quality checks
    content_warnings = tuple(check_deck_content(deck))
    error_count = count_errors(content_warnings)

    # Log all warnings
    for warning in content_warnings:
        log.warn(warning.to_log_line())

    # Write warnings.json
    warnings_json_path = None
    if log_path:
        warnings_json_path = Path(log_path).parent / "warnings.json"
        warnings_json_path.parent.mkdir(parents=True, exist_ok=True)
        warnings_data = {
            "passed": error_count == 0,
            "review_required": error_count > 0,
            "error_count": error_count,
            "warning_count": len(content_warnings) - error_count,
            "total_issues": len(content_warnings),
            "warnings": [
                {
                    "slide_id": w.slide_id,
                    "level": "blocker" if w.warning_level == "error" else w.warning_level,
                    "message": w.message,
                }
                for w in content_warnings
            ],
        }
        warnings_json_path.write_text(
            json.dumps(warnings_data, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    # Quality gate: abort if too many errors
    if error_count > max_errors:
        error_msg = f"Quality gate failed: {error_count} error(s) found, max allowed is {max_errors}. Aborting generation."
        log.error(error_msg)
        if log_path:
            final_log_path = Path(log_path)
            log.write(final_log_path)
        raise ValueError(error_msg)

    # Generate presentation
    prs = Presentation()
    prs.slide_width = Inches(theme.page.width)
    prs.slide_height = Inches(theme.page.height)

    for index, slide in enumerate(deck.slides, start=1):
        render_slide(prs, slide, theme, log, index, len(deck.slides))

    final_output = Path(output_path)
    final_output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(final_output))

    final_log_path = Path(log_path) if log_path else None
    if final_log_path is not None:
        log.write(final_log_path)

    return RenderArtifacts(
        output_path=final_output,
        log_path=final_log_path,
        warnings_path=warnings_json_path,
        slide_count=len(deck.slides),
        warnings=tuple(validated.warnings),
        content_warnings=content_warnings,
        error_count=error_count,
    )
