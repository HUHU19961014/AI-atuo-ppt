from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..theme_loader import ThemeSpec


def rgb(hex_color: str) -> RGBColor:
    hex_value = hex_color.strip().lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, theme: ThemeSpec, color_hex: str | None = None) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(theme.page.width),
        Inches(theme.page.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color_hex or theme.colors.bg)
    shape.line.fill.background()


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
    margin_left: float = 0.06,
    margin_right: float = 0.06,
    margin_top: float = 0.04,
    margin_bottom: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = Inches(margin_left)
    frame.margin_right = Inches(margin_right)
    frame.margin_top = Inches(margin_top)
    frame.margin_bottom = Inches(margin_bottom)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_hex)
    return box


def add_card(slide, left: float, top: float, width: float, height: float, theme: ThemeSpec):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(theme.colors.card_bg)
    card.line.color.rgb = rgb(theme.colors.line)
    card.line.width = Pt(1.1)
    return card


def add_bullet_list(
    slide,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
    font_size: int,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(6)
        paragraph.font.name = theme.fonts.body
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = rgb(theme.colors.text_main)
    return box


def add_page_number(slide, slide_number: int, total_slides: int, theme: ThemeSpec) -> None:
    add_textbox(
        slide,
        left=theme.page.width - theme.spacing.page_margin_right - 0.7,
        top=theme.page.height - theme.spacing.page_margin_bottom - 0.22,
        width=0.65,
        height=0.18,
        text=f"{slide_number}/{total_slides}",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.small,
        color_hex=theme.colors.text_sub,
        align=PP_ALIGN.RIGHT,
    )


def resolve_body_font_size(theme: ThemeSpec, item_count: int) -> int:
    base = theme.font_sizes.body
    if item_count <= 4:
        return base
    if item_count == 5:
        return max(theme.font_sizes.small + 2, base - 1)
    if item_count == 6:
        return max(theme.font_sizes.small + 1, base - 2)
    return max(theme.font_sizes.small, base - 3)


def add_local_image_or_placeholder(
    slide,
    *,
    image_mode: str,
    image_path: str | None,
    caption: str | None,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
    log,
) -> None:
    path = Path(image_path) if image_path else None
    if image_mode == "local_path" and path and path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
        return

    if image_mode == "local_path":
        log.warn(f"image file not found, using placeholder instead: {image_path}")
    placeholder = add_card(slide, left, top, width, height, theme)
    placeholder.fill.fore_color.rgb = rgb("#FAFAFA")
    placeholder.line.color.rgb = rgb(theme.colors.line)
    add_textbox(
        slide,
        left=left + 0.2,
        top=top + height / 2 - 0.3,
        width=width - 0.4,
        height=0.6,
        text=caption or "Image Placeholder",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.body,
        color_hex=theme.colors.text_sub,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
