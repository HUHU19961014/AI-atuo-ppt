import datetime
import gc
import re
import shutil
import subprocess
import sys
import time
import warnings
from pathlib import Path

from pptx import Presentation

from .body_renderers import apply_theme_title, fill_body_slide, fill_directory_slide
from .config import DEFAULT_MIN_TEMPLATE_SLIDES, DEFAULT_OUTPUT_DIR
from .models import BodyPageSpec
from .pipeline import plan_deck_from_html
from .powerpoint import get_powerpoint_runtime, has_powerpoint_com
from .reference_styles import build_reference_import_plan, populate_reference_body_pages
from .slide_ops import (
    clone_slide_after,
    copy_slide_xml_assets,
    ensure_last_slide,
    remove_slide,
    slide_assets_preserved,
)
from .template_manifest import TemplateManifest, load_template_manifest


def build_output_path(output_dir: Path, output_prefix: str) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    safe_prefix = re.sub(r'[<>:"/\\\\|?*]+', "_", output_prefix).strip(" ._") or "SIE_AutoPPT"
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    return output_dir / f"{safe_prefix}_{timestamp}.pptx"


def _repair_directory_slides_with_com(pptx_path: Path, source_idx: int, target_indices: list[int]) -> bool:
    if not has_powerpoint_com():
        return False
    helper = Path(__file__).resolve().parents[1] / "repair_directory_slides.py"
    command = [
        sys.executable,
        str(helper),
        str(pptx_path.resolve()),
        "--source-idx",
        str(source_idx),
        "--targets",
        *[str(target) for target in target_indices],
    ]
    for _ in range(5):
        result = subprocess.run(command, capture_output=True, text=True)
        if result.returncode == 0:
            return True
        time.sleep(1.0)
    return False


def _apply_reference_body_slides_with_com(
    pptx_path: Path,
    reference_body_path: Path | None,
    body_pages: list[BodyPageSpec],
    manifest: TemplateManifest,
) -> bool:
    import_plan = build_reference_import_plan(body_pages, manifest=manifest)
    if not import_plan:
        return False
    if not has_powerpoint_com() or reference_body_path is None or not reference_body_path.exists():
        return False
    helper = Path(__file__).resolve().parents[1] / "apply_reference_body_slides.py"
    command = [
        sys.executable,
        str(helper),
        str(pptx_path.resolve()),
        str(reference_body_path.resolve()),
        "--mapping",
        *[f"{target}={source}" for target, source in import_plan],
    ]
    last_error = ""
    for _ in range(5):
        result = subprocess.run(command, capture_output=True, text=True)
        if result.returncode == 0:
            return True
        stderr = (result.stderr or "").strip()
        stdout = (result.stdout or "").strip()
        last_error = stderr or stdout or f"exit code {result.returncode}"
        time.sleep(1.0)
    warnings.warn(f"Reference body slide import failed after retries: {last_error}", stacklevel=2)
    return False


def _render_with_preallocated_pool(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    chapter_lines: list[str],
    active_start: int,
    manifest: TemplateManifest,
):
    if not manifest.slide_pools:
        raise ValueError("Template manifest does not define slide pools.")

    used_directory_indices = list(manifest.slide_pools.directory[: len(body_pages)])
    used_body_indices = list(manifest.slide_pools.body[: len(body_pages)])
    unused_pairs = list(
        zip(
            manifest.slide_pools.directory[len(body_pages):],
            manifest.slide_pools.body[len(body_pages):],
            strict=False,
        )
    )

    for directory_idx, body_idx in sorted(unused_pairs, reverse=True):
        for slide_index in sorted((directory_idx, body_idx), reverse=True):
            if slide_index < len(prs.slides):
                remove_slide(prs, slide_index)

    directory_slides = [prs.slides[index] for index in used_directory_indices]
    body_slides = [prs.slides[index] for index in used_body_indices]

    for chapter_idx, directory_slide in enumerate(directory_slides):
        fill_directory_slide(directory_slide, chapter_lines, active_start + chapter_idx, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def _render_with_legacy_clone(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    chapter_lines: list[str],
    active_start: int,
    manifest: TemplateManifest,
):
    directory_idx = manifest.slide_roles.directory
    body_template_idx = manifest.slide_roles.body_template
    directory_slides = [prs.slides[directory_idx]]
    body_slides = [prs.slides[body_template_idx]]
    insert_after = body_template_idx
    for _ in body_pages[1:]:
        new_directory = clone_slide_after(prs, directory_idx, insert_after, keep_rel_ids=True)
        directory_slides.append(new_directory)
        insert_after += 1

        new_body = clone_slide_after(prs, body_template_idx, insert_after, keep_rel_ids=False)
        body_slides.append(new_body)
        insert_after += 1

    for chapter_idx, directory_slide in enumerate(directory_slides):
        fill_directory_slide(directory_slide, chapter_lines, active_start + chapter_idx, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def _refresh_legacy_directory_clones(
    pptx_path: Path,
    chapter_lines: list[str],
    active_start: int,
    body_page_count: int,
    manifest: TemplateManifest,
):
    directory_idx = manifest.slide_roles.directory
    targets = [directory_idx + 1 + i * 2 for i in range(1, body_page_count)]
    if not targets:
        return True

    source_idx = directory_idx + 1
    for strategy in ("zip", "com"):
        for _ in range(3):
            if strategy == "zip":
                if not copy_slide_xml_assets(pptx_path, source_idx=source_idx, target_indices=targets):
                    continue
            else:
                if not _repair_directory_slides_with_com(pptx_path, source_idx=source_idx, target_indices=targets):
                    continue

            prs_reloaded = Presentation(str(pptx_path))
            fill_directory_slide(prs_reloaded.slides[directory_idx], chapter_lines, active_start, manifest)
            for offset, directory_slide_no in enumerate(targets, start=1):
                slide_index = directory_slide_no - 1
                if slide_index < len(prs_reloaded.slides):
                    fill_directory_slide(prs_reloaded.slides[slide_index], chapter_lines, active_start + offset, manifest)
            prs_reloaded.save(str(pptx_path))
            prs_reloaded = None
            gc.collect()

            if slide_assets_preserved(pptx_path, source_idx=source_idx, target_indices=targets):
                return True
            time.sleep(1.0)

    return False


def _warn_if_reference_import_disabled(body_pages: list[BodyPageSpec], reference_body_path: Path | None):
    if not any(page.reference_style_id for page in body_pages):
        return
    if reference_body_path is None or not reference_body_path.exists():
        warnings.warn(
            "Reference style library is unavailable; using native fallback renderers for reference-style pages.",
            stacklevel=2,
        )
        return
    if not has_powerpoint_com():
        runtime = get_powerpoint_runtime()
        warnings.warn(
            f"{runtime.reason} Native fallback renderers will be used for reference-style pages.",
            stacklevel=2,
        )


def _refresh_preallocated_directory_assets(pptx_path: Path, body_page_count: int, manifest: TemplateManifest) -> bool:
    if not manifest.slide_pools:
        return True
    target_indices = [index + 1 for index in manifest.slide_pools.directory[1:body_page_count]]
    if not target_indices:
        return True
    source_idx = manifest.slide_pools.directory[0] + 1
    if not copy_slide_xml_assets(pptx_path, source_idx=source_idx, target_indices=target_indices):
        return False
    return slide_assets_preserved(pptx_path, source_idx=source_idx, target_indices=target_indices)


def generate_ppt(
    template_path: Path,
    html_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    chapters: int,
    active_start: int,
    output_dir: Path | None = None,
):
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    if not html_path.exists():
        raise FileNotFoundError(f"HTML not found: {html_path}")

    manifest = load_template_manifest(template_path=template_path)
    deck_plan = plan_deck_from_html(html_path, chapters)
    deck = deck_plan.deck
    body_pages = deck.body_pages
    chapter_lines = deck_plan.chapter_lines
    pattern_ids = deck_plan.pattern_ids

    final_output_dir = output_dir or DEFAULT_OUTPUT_DIR
    out = build_output_path(final_output_dir, output_prefix)
    shutil.copy2(template_path, out)

    prs = Presentation(str(out))
    if len(prs.slides) < DEFAULT_MIN_TEMPLATE_SLIDES:
        raise ValueError(
            f"\u6a21\u677f\u9875\u6570\u4e0d\u8db3\uff0c\u81f3\u5c11\u9700\u8981 {DEFAULT_MIN_TEMPLATE_SLIDES} \u9875\uff0c\u5b9e\u9645\u4e3a {len(prs.slides)} \u9875\u3002"
        )

    apply_theme_title(prs, deck.cover_title, manifest)
    if manifest.slide_pools and manifest.slide_pools.ending is not None and manifest.slide_pools.ending < len(prs.slides):
        thanks_slide_id = int(prs.slides._sldIdLst[manifest.slide_pools.ending].id)
    else:
        thanks_slide_id = int(prs.slides._sldIdLst[len(prs.slides) - 1].id)

    used_preallocated_pool = manifest.supports_preallocated_pool(len(body_pages), len(prs.slides))
    if used_preallocated_pool:
        _render_with_preallocated_pool(prs, body_pages, chapter_lines, active_start, manifest)
    else:
        warnings.warn(
            "Template does not provide a preallocated slide pool; falling back to runtime cloning for this deck.",
            stacklevel=2,
        )
        _render_with_legacy_clone(prs, body_pages, chapter_lines, active_start, manifest)

    ensure_last_slide(prs, thanks_slide_id)
    prs.save(str(out))
    prs = None
    gc.collect()

    _warn_if_reference_import_disabled(body_pages, reference_body_path)
    reference_import_applied = _apply_reference_body_slides_with_com(out, reference_body_path, body_pages, manifest)
    if reference_import_applied:
        populate_reference_body_pages(out, body_pages, manifest=manifest)

    if used_preallocated_pool and not _refresh_preallocated_directory_assets(out, len(body_pages), manifest):
        warnings.warn(
            "Preallocated directory slide assets could not be fully refreshed after save; manual review is recommended.",
            stacklevel=2,
        )

    if not used_preallocated_pool:
        if not _refresh_legacy_directory_clones(out, chapter_lines, active_start, len(body_pages), manifest):
            warnings.warn(
                "Directory slide clone repair did not fully preserve template image assets; keeping generated deck and surfacing the risk in QA.",
                stacklevel=2,
            )

    return out, pattern_ids, chapter_lines
